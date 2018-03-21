# coding:utf-8
"""
python写pptx，demo
"""
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

demo_pptx_path = r'C:\Users\stephen\Desktop\1.pptx'

prs = Presentation(demo_pptx_path)
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = ChartData()
chart_data.categories = ['east', 'west', 'midwest']
chart_data.add_series('series1', (19.2, 21.4, 16.7))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)
chart = graphic_frame.chart
# chart.font.size = Pt(13)

prs.save(r'C:\Users\stephen\Desktop\chart-01.pptx')
